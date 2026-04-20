import io
import os
import uuid
from dataclasses import dataclass, asdict
from typing import List, Literal, Optional

import streamlit as st
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE_PATH = "templates.pptx"
SLIDE_W = 13.333  # inches for 16:9
SLIDE_H = 7.5
PREVIEW_W = 1280
PREVIEW_H = 720

TITLE_GENERAL_STYLE = {
    "x": 1.0,
    "y": 2.35,
    "w": 11.3,
    "h": 1.0,
    "font_size": 28,
    "bold": True,
    "align": "center",
}

TITLE_SECTION_STYLE = {
    "x": 1.0,
    "y": 2.35,
    "w": 11.3,
    "h": 1.0,
    "font_size": 24,
    "bold": True,
    "align": "center",
}

CONTENT_HEADER_STYLE = {
    "x": 0.7,
    "y": 0.18,
    "w": 11.9,
    "h": 0.5,
    "font_size": 22,
    "bold": True,
    "align": "center",
}

CONTENT_ZONE = {
    "x": 0.65,
    "y": 1.05,
    "w": 12.0,
    "h": 6.0,
}


@dataclass
class ContentItem:
    kind: Literal["text", "image"]
    x: float
    y: float
    w: float
    h: float
    text: str = ""
    font_size: int = 14
    bold: bool = False
    image_bytes: Optional[bytes] = None
    filename: str = ""


@dataclass
class SlideSpec:
    slide_type: Literal["title_general", "title_section", "content"]
    title: str
    items: Optional[List[ContentItem]] = None


# ---------- Helpers ----------

def inches_to_px_x(value: float) -> int:
    return int(value / SLIDE_W * PREVIEW_W)


def inches_to_px_y(value: float) -> int:
    return int(value / SLIDE_H * PREVIEW_H)


def clamp_to_content_zone(x: float, y: float, w: float, h: float):
    min_x = CONTENT_ZONE["x"]
    min_y = CONTENT_ZONE["y"]
    max_x = CONTENT_ZONE["x"] + CONTENT_ZONE["w"]
    max_y = CONTENT_ZONE["y"] + CONTENT_ZONE["h"]

    x = max(min_x, min(x, max_x - 0.2))
    y = max(min_y, min(y, max_y - 0.2))
    w = min(w, max_x - x)
    h = min(h, max_y - y)
    return round(x, 2), round(y, 2), round(w, 2), round(h, 2)


def get_font(size: int, bold: bool = False):
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            return ImageFont.truetype(path, size=size)
    return ImageFont.load_default()


def add_wrapped_text(draw: ImageDraw.ImageDraw, box, text: str, font, fill=(20, 20, 20), align="left"):
    x, y, w, h = box
    words = text.split()
    if not words:
        return

    lines = []
    current = words[0]
    for word in words[1:]:
        trial = current + " " + word
        trial_bbox = draw.textbbox((0, 0), trial, font=font)
        if trial_bbox[2] - trial_bbox[0] <= w:
            current = trial
        else:
            lines.append(current)
            current = word
    lines.append(current)

    line_heights = []
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font)
        line_heights.append(bbox[3] - bbox[1])
    total_h = sum(line_heights) + max(0, len(lines) - 1) * 6

    current_y = y + max(0, (h - total_h) // 2)
    for idx, line in enumerate(lines):
        bbox = draw.textbbox((0, 0), line, font=font)
        line_w = bbox[2] - bbox[0]
        line_h = bbox[3] - bbox[1]
        if align == "center":
            text_x = x + max(0, (w - line_w) // 2)
        else:
            text_x = x
        draw.text((text_x, current_y), line, font=font, fill=fill)
        current_y += line_h + 6


def render_slide_preview(slide_spec: SlideSpec) -> Image.Image:
    img = Image.new("RGB", (PREVIEW_W, PREVIEW_H), "white")
    draw = ImageDraw.Draw(img)

    # Simple preview background approximating the template logic.
    if slide_spec.slide_type == "title_general":
        draw.rectangle([0, 0, PREVIEW_W, PREVIEW_H], fill=(42, 73, 108))
        box = (
            inches_to_px_x(TITLE_GENERAL_STYLE["x"]),
            inches_to_px_y(TITLE_GENERAL_STYLE["y"]),
            inches_to_px_x(TITLE_GENERAL_STYLE["w"]),
            inches_to_px_y(TITLE_GENERAL_STYLE["h"]),
        )
        font = get_font(34, True)
        add_wrapped_text(draw, box, slide_spec.title, font, fill=(255, 255, 255), align="center")

    elif slide_spec.slide_type == "title_section":
        draw.rectangle([0, 0, PREVIEW_W, PREVIEW_H], fill=(94, 131, 164))
        box = (
            inches_to_px_x(TITLE_SECTION_STYLE["x"]),
            inches_to_px_y(TITLE_SECTION_STYLE["y"]),
            inches_to_px_x(TITLE_SECTION_STYLE["w"]),
            inches_to_px_y(TITLE_SECTION_STYLE["h"]),
        )
        font = get_font(30, True)
        add_wrapped_text(draw, box, slide_spec.title, font, fill=(255, 255, 255), align="center")

    else:
        header_h = inches_to_px_y(0.8)
        draw.rectangle([0, 0, PREVIEW_W, header_h], fill=(42, 73, 108))
        cz = (
            inches_to_px_x(CONTENT_ZONE["x"]),
            inches_to_px_y(CONTENT_ZONE["y"]),
            inches_to_px_x(CONTENT_ZONE["x"] + CONTENT_ZONE["w"]),
            inches_to_px_y(CONTENT_ZONE["y"] + CONTENT_ZONE["h"]),
        )
        draw.rounded_rectangle(cz, outline=(210, 210, 210), width=3, radius=10)

        title_box = (
            inches_to_px_x(CONTENT_HEADER_STYLE["x"]),
            inches_to_px_y(CONTENT_HEADER_STYLE["y"]),
            inches_to_px_x(CONTENT_HEADER_STYLE["w"]),
            inches_to_px_y(CONTENT_HEADER_STYLE["h"]),
        )
        font = get_font(26, True)
        add_wrapped_text(draw, title_box, slide_spec.title, font, fill=(255, 255, 255), align="center")

        for item in slide_spec.items or []:
            x = inches_to_px_x(item.x)
            y = inches_to_px_y(item.y)
            w = inches_to_px_x(item.w)
            h = inches_to_px_y(item.h)
            rect = [x, y, x + w, y + h]

            if item.kind == "text":
                draw.rounded_rectangle(rect, outline=(120, 120, 120), width=2, radius=8)
                font = get_font(max(14, int(item.font_size * 1.5)), item.bold)
                add_wrapped_text(draw, (x + 12, y + 12, w - 24, h - 24), item.text or "Texte", font)
            else:
                draw.rounded_rectangle(rect, outline=(120, 120, 120), width=2, radius=8, fill=(245, 245, 245))
                if item.image_bytes:
                    try:
                        im = Image.open(io.BytesIO(item.image_bytes)).convert("RGB")
                        im.thumbnail((max(40, w - 8), max(40, h - 8)))
                        paste_x = x + (w - im.width) // 2
                        paste_y = y + (h - im.height) // 2
                        img.paste(im, (paste_x, paste_y))
                    except Exception:
                        draw.line([x + 10, y + 10, x + w - 10, y + h - 10], fill=(160, 160, 160), width=3)
                        draw.line([x + w - 10, y + 10, x + 10, y + h - 10], fill=(160, 160, 160), width=3)
                else:
                    draw.line([x + 10, y + 10, x + w - 10, y + h - 10], fill=(160, 160, 160), width=3)
                    draw.line([x + w - 10, y + 10, x + 10, y + h - 10], fill=(160, 160, 160), width=3)
                    ph_font = get_font(20, True)
                    add_wrapped_text(draw, (x + 12, y + 12, w - 24, h - 24), "Image", ph_font, align="center")

    return img


def add_textbox(slide, cfg, text: str):
    box = slide.shapes.add_textbox(
        Inches(cfg["x"]), Inches(cfg["y"]), Inches(cfg["w"]), Inches(cfg["h"])
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER if cfg["align"] == "center" else PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(cfg["font_size"])
    run.font.bold = cfg["bold"]
    run.font.name = "Aptos"


def build_ppt(slides: List[SlideSpec]) -> bytes:
    prs = Presentation(TEMPLATE_PATH)

    # keep only template slides already in file and append generated slides after them
    for spec in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        if spec.slide_type == "title_general":
            add_textbox(slide, TITLE_GENERAL_STYLE, spec.title)
        elif spec.slide_type == "title_section":
            add_textbox(slide, TITLE_SECTION_STYLE, spec.title)
        else:
            add_textbox(slide, CONTENT_HEADER_STYLE, spec.title)
            for item in spec.items or []:
                if item.kind == "text":
                    shape = slide.shapes.add_textbox(
                        Inches(item.x), Inches(item.y), Inches(item.w), Inches(item.h)
                    )
                    tf = shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = item.text
                    p.alignment = PP_ALIGN.LEFT
                    run = p.runs[0]
                    run.font.size = Pt(item.font_size)
                    run.font.bold = item.bold
                    run.font.name = "Aptos"
                elif item.kind == "image" and item.image_bytes:
                    slide.shapes.add_picture(
                        io.BytesIO(item.image_bytes),
                        Inches(item.x), Inches(item.y),
                        width=Inches(item.w), height=Inches(item.h)
                    )

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# ---------- UI ----------
st.set_page_config(page_title="Générateur Solimed", layout="wide")
st.title("Générateur de présentation Solimed")
st.caption("Version simple avec preview visuel et export PowerPoint")

if "slides" not in st.session_state:
    st.session_state.slides = []

with st.sidebar:
    st.header("Nouvelle slide")
    slide_type_ui = st.selectbox(
        "Type de slide",
        ["Titre général", "Titre intermédiaire", "Contenu"],
    )
    type_map = {
        "Titre général": "title_general",
        "Titre intermédiaire": "title_section",
        "Contenu": "content",
    }
    slide_type = type_map[slide_type_ui]
    title = st.text_input("Titre")

    items: List[ContentItem] = []
    if slide_type == "content":
        st.markdown("### Éléments dans la zone blanche")
        n_items = st.number_input("Nombre d'éléments", min_value=0, max_value=8, value=1, step=1)

        for i in range(n_items):
            st.markdown(f"**Élément {i + 1}**")
            kind_ui = st.selectbox("Type", ["Texte", "Image"], key=f"kind_{i}")
            kind = "text" if kind_ui == "Texte" else "image"

            col1, col2 = st.columns(2)
            with col1:
                x = st.number_input("x", 0.0, 12.5, 1.0, 0.1, key=f"x_{i}")
                w = st.number_input("largeur", 0.3, 12.0, 4.0, 0.1, key=f"w_{i}")
            with col2:
                y = st.number_input("y", 0.8, 7.0, 1.4, 0.1, key=f"y_{i}")
                h = st.number_input("hauteur", 0.3, 6.0, 2.0, 0.1, key=f"h_{i}")

            x, y, w, h = clamp_to_content_zone(x, y, w, h)

            if kind == "text":
                text = st.text_area("Texte", key=f"text_{i}")
                font_size = st.slider("Taille police", 10, 28, 14, key=f"font_{i}")
                bold = st.checkbox("Gras", key=f"bold_{i}")
                items.append(ContentItem(kind="text", x=x, y=y, w=w, h=h, text=text, font_size=font_size, bold=bold))
            else:
                image_file = st.file_uploader("Image", type=["png", "jpg", "jpeg"], key=f"img_{i}")
                image_bytes = image_file.getvalue() if image_file else None
                filename = image_file.name if image_file else ""
                items.append(ContentItem(kind="image", x=x, y=y, w=w, h=h, image_bytes=image_bytes, filename=filename))

    current_slide = SlideSpec(slide_type=slide_type, title=title, items=items if slide_type == "content" else None)

    if st.button("Ajouter la slide", use_container_width=True):
        if not title.strip():
            st.error("Le titre est obligatoire.")
        else:
            st.session_state.slides.append(current_slide)
            st.success("Slide ajoutée.")

    if st.button("Supprimer la dernière", use_container_width=True):
        if st.session_state.slides:
            st.session_state.slides.pop()
            st.success("Dernière slide supprimée.")

    if st.button("Vider la présentation", use_container_width=True):
        st.session_state.slides = []
        st.success("Présentation vidée.")

# Main area
left, right = st.columns([1.1, 1.4])

with left:
    st.subheader("Preview de la slide en cours")
    preview = render_slide_preview(current_slide)
    st.image(preview, use_container_width=True)

with right:
    st.subheader("Slides de la présentation")
    if not st.session_state.slides:
        st.info("Aucune slide ajoutée pour le moment.")
    else:
        for idx, slide in enumerate(st.session_state.slides, start=1):
            with st.expander(f"Slide {idx} — {slide.title}", expanded=False):
                st.write(f"Type : {slide.slide_type}")
                st.image(render_slide_preview(slide), use_container_width=True)

st.divider()
st.subheader("Export")

if not os.path.exists(TEMPLATE_PATH):
    st.error("Le fichier templates.pptx doit être à la racine du projet.")
else:
    if st.session_state.slides:
        ppt_bytes = build_ppt(st.session_state.slides)
        st.download_button(
            "Télécharger le PPT",
            data=ppt_bytes,
            file_name=f"presentation_{uuid.uuid4().hex[:8]}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )
    else:
        st.warning("Ajoute au moins une slide pour générer le PPT.")
